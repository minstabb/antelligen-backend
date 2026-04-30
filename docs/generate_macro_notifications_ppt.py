"""Generate macro-and-schedule-notifications.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Palette (matches existing decks)
DARK_BG   = RGBColor(0x0D, 0x2B, 0x55)
HEADER_BG = RGBColor(0x3A, 0x59, 0xD1)
ACCENT    = RGBColor(0x5B, 0xC4, 0xFF)
GREEN     = RGBColor(0xA6, 0xE3, 0xA1)
YELLOW    = RGBColor(0xF9, 0xE2, 0xAF)
RED       = RGBColor(0xF3, 0x8B, 0xA8)
PURPLE    = RGBColor(0xCB, 0xA6, 0xF7)
TEAL      = RGBColor(0x94, 0xE2, 0xD8)
ORANGE    = RGBColor(0xFA, 0xB3, 0x87)
SUBTEXT   = RGBColor(0xB8, 0xD4, 0xF0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG   = RGBColor(0x16, 0x3E, 0x75)
CARD_BG2  = RGBColor(0x1A, 0x48, 0x85)
BORDER    = RGBColor(0x2D, 0x6A, 0xB0)

blank = prs.slide_layouts[6]


def set_bg(slide, color=DARK_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, left, top, w, h, fill, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, left, top, w, h, size=11, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def multitxt(slide, lines, left, top, w, h, size=10, color=WHITE, bullet="• "):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = (bullet + line) if bullet else line
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.85, HEADER_BG)
    txt(slide, title, 0.4, 0.15, 12.5, 0.55, size=24, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.55, 12.5, 0.3, size=11, color=SUBTEXT)


def card(slide, left, top, w, h, title, lines, title_color=ACCENT, body_size=10,
         bg=CARD_BG, line=BORDER):
    rect(slide, left, top, w, h, bg, line=line, line_w=1)
    txt(slide, title, left + 0.15, top + 0.1, w - 0.3, 0.4, size=13, bold=True, color=title_color)
    multitxt(slide, lines, left + 0.15, top + 0.55, w - 0.3, h - 0.6, size=body_size)


# ══════════════════════════════════════════════════════════
# Slide 1: Title
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
rect(s, 0, 2.4, 13.33, 2.7, CARD_BG)
txt(s, "거시 경제 현황판 + 경제일정 알림",
    0, 2.55, 13.33, 0.9, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Macro Dashboard  ·  Schedule Notifications",
    0, 3.45, 13.33, 0.6, size=22, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Antelligen Backend — 두 핵심 사용자 기능 요약",
    0, 4.15, 13.33, 0.5, size=15, color=SUBTEXT, align=PP_ALIGN.CENTER)
txt(s, "2026.04 · Antelligen AI",
    0, 6.7, 13.33, 0.4, size=12, color=SUBTEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# Slide 2: 한 눈에 (역할 분담 + 공유 데이터)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "두 기능 한 눈에", "역할 분담과 공유 데이터")

# Macro card
rect(s, 0.4, 1.1, 6.0, 3.2, CARD_BG, line=ACCENT, line_w=2)
rect(s, 0.4, 1.1, 6.0, 0.55, ACCENT)
txt(s, "거시 경제 현황판  /  macro 도메인", 0.55, 1.18, 5.7, 0.4,
    size=15, bold=True, color=DARK_BG)
multitxt(s, [
    "Risk-on / Risk-off 즉시 응답",
    "일 1회 LLM 갱신 (01:00 KST)",
    "메모리 + Redis 25h 캐시",
    "학습 노트 + YouTube + 월가 IB 페르소나",
    "contextual + baseline 듀얼 판단",
], 0.55, 1.8, 5.7, 2.4, size=12, bullet="• ")

# Notifications card
rect(s, 6.95, 1.1, 6.0, 3.2, CARD_BG, line=ORANGE, line_w=2)
rect(s, 6.95, 1.1, 6.0, 0.55, ORANGE)
txt(s, "경제일정 알림  /  schedule 도메인", 7.1, 1.18, 5.7, 0.4,
    size=15, bold=True, color=DARK_BG)
multitxt(s, [
    "LLM 영향 분석 저장 시 알림 발사",
    "DB 기록 + SSE 실시간 푸시",
    "읽음 상태 (단건 / 전체) 관리",
    "FOMC collapse + (M/D) suffix 라벨",
    "분석 입력에 매크로 지표 13종 주입",
], 7.1, 1.8, 5.7, 2.4, size=12, bullet="• ")

# Shared data ribbon
rect(s, 0.4, 4.5, 12.55, 0.55, HEADER_BG)
txt(s, "공유 데이터: 매크로 지표 스냅샷 13종 (금리·유가·환율·VIX·DXY·지수·금)",
    0.4, 4.6, 12.55, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom value props
props = [
    ("일관된 시장 톤", GREEN, "거시 현황의 contextual 판단과\n경제일정 영향 분석이 같은 데이터로 강화"),
    ("빠른 응답", YELLOW, "LLM 호출은 백그라운드/캐시\n사용자 응답은 캐시 hit"),
    ("실시간 알림", PURPLE, "SSE 스트림으로 분석 완료 즉시\n프론트 종모양 UI 갱신"),
]
for i, (title, color, body) in enumerate(props):
    left = 0.4 + i * 4.21
    rect(s, left, 5.25, 4.05, 1.85, CARD_BG2, line=color, line_w=2)
    rect(s, left, 5.25, 4.05, 0.4, color)
    txt(s, title, left + 0.15, 5.3, 3.75, 0.35, size=12, bold=True, color=DARK_BG)
    multitxt(s, body.split("\n"), left + 0.2, 5.7, 3.7, 1.3, size=11, bullet="• ")


# ══════════════════════════════════════════════════════════
# Slide 3: Macro Dashboard — 흐름 + 데이터 소스
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "거시 경제 현황판 — 데이터 흐름", "APScheduler → LLM 듀얼 판단 → 캐시 → 즉시 응답")

# pipeline boxes
nodes = [
    ("APScheduler\n01:00 daily", ACCENT),
    ("Judge\nUseCase", GREEN),
    ("Sources\n(노트·YT·LLM)", YELLOW),
    ("Snapshot\nStore", ORANGE),
    ("Memory +\nRedis 25h", PURPLE),
    ("GET /macro/\nmarket-risk", TEAL),
]
top = 1.3
node_w = 1.95
gap = 0.2
x = 0.4
for label, c in nodes:
    rect(s, x, top, node_w, 1.1, c)
    tb = s.shapes.add_textbox(Inches(x), Inches(top + 0.15), Inches(node_w), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = DARK_BG
    x += node_w + gap

# arrow line
rect(s, 0.4, 2.5, 12.5, 0.05, ACCENT)

# 3 source cards
sources = [
    ("학습 노트", GREEN, [
        "StudyNoteFileReader",
        "로컬 파일 시스템",
        "contextual 판단 근거",
    ]),
    ("YouTube 영상", YELLOW, [
        "Antelligen 채널 최근 7일",
        "youtube_macro_video_client.py",
        "reference_videos 4건 노출",
    ]),
    ("OpenAI GPT", ORANGE, [
        "langchain_risk_judgement_adapter.py",
        "월가 IB 페르소나 (GS · MS · JPM)",
        "contextual + baseline 듀얼 출력",
    ]),
]
top2 = 2.85
for i, (title, color, items) in enumerate(sources):
    left = 0.4 + i * 4.21
    rect(s, left, top2, 4.05, 2.4, CARD_BG, line=color, line_w=2)
    rect(s, left, top2, 4.05, 0.5, color)
    txt(s, title, left + 0.15, top2 + 0.07, 3.75, 0.4, size=13, bold=True, color=DARK_BG)
    multitxt(s, items, left + 0.2, top2 + 0.6, 3.7, 1.7, size=11, bullet="• ")

# Caching note
rect(s, 0.4, 5.5, 12.55, 1.6, CARD_BG2, line=ACCENT, line_w=1)
txt(s, "캐시 전략 (Hot reload 안전)", 0.6, 5.6, 12.0, 0.4, size=14, bold=True, color=ACCENT)
multitxt(s, [
    "메모리: MarketRiskSnapshotStore (스레드-세이프, 프로세스 싱글톤)",
    "Redis: TTL 25h — 프로세스 재시작 시 직전 스냅샷 복원 → YouTube/LLM 재호출 회피",
    "사용자 응답은 캐시 hit 만 — LLM 호출은 매일 01:00 한 번뿐",
], 0.6, 5.95, 12.0, 1.1, size=11, bullet="✓ ")


# ══════════════════════════════════════════════════════════
# Slide 4: Macro — 듀얼 판단 + 응답 구조
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "거시 현황판 — 듀얼 판단 & 응답 스키마", "contextual + baseline / 출처 표기 일원화")

# Left: dual judgement diagram
rect(s, 0.4, 1.1, 6.4, 3.0, CARD_BG, line=ACCENT, line_w=2)
txt(s, "이중 판단 체계", 0.6, 1.2, 6.0, 0.4, size=14, bold=True, color=ACCENT)

rect(s, 0.6, 1.75, 2.9, 2.25, GREEN)
txt(s, "Contextual", 0.6, 1.85, 2.9, 0.4, size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
multitxt(s, [
    "학습 노트 + YouTube",
    "RISK_ON / RISK_OFF / UNKNOWN",
    "근거 3줄",
    "프로젝트 컨텍스트 반영",
], 0.7, 2.3, 2.7, 1.6, size=10, color=DARK_BG, bullet="· ")

rect(s, 3.7, 1.75, 2.9, 2.25, ORANGE)
txt(s, "Baseline", 3.7, 1.85, 2.9, 0.4, size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
multitxt(s, [
    "월가 IB 페르소나",
    "GS · MS · JPM Research 톤",
    "근거 3줄",
    "일반 시장 컨센서스 판단",
], 3.8, 2.3, 2.7, 1.6, size=10, color=DARK_BG, bullet="· ")

# Right: response schema
rect(s, 7.0, 1.1, 5.95, 5.7, CARD_BG, line=BORDER, line_w=1)
txt(s, "MarketRiskJudgementResponse", 7.15, 1.2, 5.7, 0.4, size=14, bold=True, color=YELLOW)

rows = [
    ("status", "최종 판단 (RISK_ON/OFF/UNKNOWN)"),
    ("contextual_status", "학습/영상 기반 판단"),
    ("contextual_reasons[3]", "3줄 근거"),
    ("baseline_status", "월가 페르소나 판단"),
    ("baseline_reasons[3]", "3줄 근거"),
    ("reference_videos[]", "id · title · published_at · url"),
    ("note_available", "학습 노트 보유 여부"),
    ("updated_at", "스냅샷 갱신 시각"),
]
top = 1.7
for field, desc in rows:
    rect(s, 7.15, top, 5.65, 0.55, CARD_BG2)
    txt(s, field, 7.3, top + 0.1, 2.3, 0.35, size=11, bold=True, color=ACCENT)
    txt(s, desc, 9.6, top + 0.1, 3.1, 0.35, size=10, color=WHITE)
    top += 0.62

# Bottom: 출처 표기
rect(s, 0.4, 4.3, 6.4, 2.8, CARD_BG2, line=YELLOW, line_w=2)
txt(s, "출처 표기 일원화", 0.6, 4.4, 6.0, 0.4, size=14, bold=True, color=YELLOW)
multitxt(s, [
    "모든 응답을 'Antelligen AI 자체 분석'으로 표기",
    "유튜브 채널명 · 영상명 · 외부 리서치 기관명 노출 금지",
    "월가 IB 페르소나 · 한국어 존댓말로 일관 응답",
    "프롬프트 규칙: langchain_risk_judgement_adapter.py:37-40",
], 0.6, 4.85, 6.0, 2.2, size=11, bullet="✓ ")


# ══════════════════════════════════════════════════════════
# Slide 5: Notifications — 트리거 + 흐름
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "경제일정 알림 — 발생 흐름", "분석 저장 → DB 기록 + SSE 푸시")

# trigger row
rect(s, 0.4, 1.1, 12.55, 0.6, CARD_BG2, line=ACCENT, line_w=1)
txt(s, "트리거", 0.55, 1.18, 1.5, 0.4, size=13, bold=True, color=ACCENT)
txt(s, "POST /schedule/event-analysis/run     ·     GET /schedule/event-analysis (lazy)",
    2.0, 1.22, 10.5, 0.4, size=12, color=WHITE)

# pipeline
nodes2 = [
    ("RunEvent\nAnalysis\nUseCase", ACCENT),
    ("매크로 지표\n스냅샷 13종\n(FRED+Yahoo)", GREEN),
    ("OpenAI\nEvent Impact\nAnalyzer", YELLOW),
    ("Schedule\nNotification\nPublisher", ORANGE),
    ("DB INSERT\n+\nSSE Broadcast", PURPLE),
]
top = 2.0
node_w = 2.3
gap = 0.2
x = 0.6
for label, c in nodes2:
    rect(s, x, top, node_w, 1.5, c)
    tb = s.shapes.add_textbox(Inches(x), Inches(top + 0.15), Inches(node_w), Inches(1.25))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = DARK_BG
    x += node_w + gap

# arrows
rect(s, 0.6, 3.7, 12.1, 0.05, ACCENT)

# 2 result lines
rect(s, 0.4, 4.0, 12.55, 1.0, CARD_BG, line=GREEN, line_w=2)
rect(s, 0.4, 4.0, 0.18, 1.0, GREEN)
txt(s, "DB INSERT", 0.6, 4.1, 2.0, 0.4, size=13, bold=True, color=GREEN)
txt(s, "schedule_notifications row 영구 보관 → GET /schedule/notifications 로 목록 조회 가능",
    2.7, 4.2, 9.8, 0.7, size=12, color=WHITE)

rect(s, 0.4, 5.15, 12.55, 1.0, CARD_BG, line=PURPLE, line_w=2)
rect(s, 0.4, 5.15, 0.18, 1.0, PURPLE)
txt(s, "SSE Broadcast", 0.6, 5.25, 2.5, 0.4, size=13, bold=True, color=PURPLE)
txt(s, "asyncio.Queue 기반 pub/sub → GET /schedule/notifications/stream 구독자에 즉시 fan-out",
    2.7, 5.35, 9.8, 0.7, size=12, color=WHITE)

# Note: separation
rect(s, 0.4, 6.3, 12.55, 0.85, CARD_BG2, line=YELLOW, line_w=1)
txt(s, "✓ 저장과 푸시는 분리 — DB INSERT 실패해도 SSE 브로드캐스트는 최선 노력으로 수행",
    0.55, 6.5, 12.3, 0.5, size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# Slide 6: Notifications — 저장 스키마 + API
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "경제일정 알림 — 스키마 & API", "schedule_notifications 테이블과 4종 엔드포인트")

# Left: table schema
rect(s, 0.4, 1.1, 6.4, 5.85, CARD_BG, line=BORDER, line_w=1)
txt(s, "schedule_notifications", 0.55, 1.2, 6.1, 0.4, size=14, bold=True, color=ACCENT)
txt(s, "schedule_notification_orm.py:9-28", 0.55, 1.55, 6.1, 0.3, size=9, color=SUBTEXT, italic=True)

cols = [
    ("id", "PK / autoincrement", ACCENT),
    ("event_id", "FK → economic_events.id", GREEN),
    ("event_title", "VARCHAR(255) — 일정 제목", WHITE),
    ("analysis_id", "FK nullable", WHITE),
    ("success", "BOOLEAN", WHITE),
    ("stored_at", "DATETIME — 분석 저장 시각", WHITE),
    ("error_message", "TEXT — 실패 메시지", RED),
    ("read_at", "DATETIME nullable — 읽음 시각", YELLOW),
    ("created_at", "DATETIME — 알림 생성 시각", WHITE),
]
top = 2.0
for col, desc, c in cols:
    rect(s, 0.55, top, 6.1, 0.45, CARD_BG2)
    txt(s, col, 0.65, top + 0.07, 1.8, 0.35, size=11, bold=True, color=c)
    txt(s, desc, 2.55, top + 0.08, 4.0, 0.35, size=10, color=WHITE)
    top += 0.5

txt(s, "INDEX  (read_at, created_at) · (event_id)",
    0.65, 6.55, 6.0, 0.35, size=10, italic=True, color=SUBTEXT)

# Right: API endpoints
rect(s, 7.0, 1.1, 5.95, 5.85, CARD_BG, line=BORDER, line_w=1)
txt(s, "API 엔드포인트", 7.15, 1.2, 5.7, 0.4, size=14, bold=True, color=ORANGE)

apis = [
    ("GET", "/schedule/notifications", "목록 (limit, unread_only)", GREEN),
    ("POST", "/schedule/notifications/{id}/read", "개별 읽음", YELLOW),
    ("POST", "/schedule/notifications/read-all", "전체 읽음", YELLOW),
    ("GET", "/schedule/notifications/stream", "SSE 실시간 구독", PURPLE),
]
top = 1.7
for method, path, desc, c in apis:
    rect(s, 7.15, top, 5.65, 1.15, CARD_BG2, line=c, line_w=2)
    rect(s, 7.15, top, 0.7, 1.15, c)
    txt(s, method, 7.15, top + 0.4, 0.7, 0.35, size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, path, 7.95, top + 0.1, 4.7, 0.4, size=11, bold=True, color=WHITE)
    txt(s, desc, 7.95, top + 0.55, 4.7, 0.5, size=10, color=SUBTEXT)
    top += 1.27


# ══════════════════════════════════════════════════════════
# Slide 7: Notifications — SSE Broadcaster 구조
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "SSE 실시간 푸시 — Broadcaster 구조",
       "asyncio.Queue 기반 인-프로세스 pub/sub")

# Publisher
rect(s, 0.5, 1.3, 3.3, 1.6, ORANGE)
txt(s, "Publisher", 0.5, 1.5, 3.3, 0.4, size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
txt(s, "publish(payload)", 0.5, 1.95, 3.3, 0.35, size=11, color=DARK_BG, align=PP_ALIGN.CENTER)
txt(s, "RunEventAnalysis 끝 시점", 0.5, 2.35, 3.3, 0.35, size=10, italic=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Broadcaster (singleton)
rect(s, 4.3, 1.0, 4.7, 2.3, HEADER_BG, line=ACCENT, line_w=3)
txt(s, "NotificationBroadcaster", 4.3, 1.15, 4.7, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "(싱글톤)", 4.3, 1.55, 4.7, 0.3, size=10, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
multitxt(s, [
    "asyncio.Queue 리스트 보관",
    "publish() = fan-out to all queues",
    "subscribe() = 새 큐 등록",
], 4.45, 1.95, 4.4, 1.3, size=11, color=WHITE, bullet="• ")

# Subscribers (3 fan-out)
subs = [
    ("Subscriber A", GREEN, 9.5, 1.1),
    ("Subscriber B", GREEN, 9.5, 2.1),
    ("Subscriber C", GREEN, 9.5, 3.1),
]
for label, c, x, y in subs:
    rect(s, x, y, 3.4, 0.85, c)
    txt(s, label, x, y + 0.25, 3.4, 0.4, size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# arrows (rectangles as lines)
rect(s, 3.85, 2.1, 0.4, 0.06, WHITE)  # publisher → broadcaster
rect(s, 9.05, 1.55, 0.4, 0.06, WHITE)
rect(s, 9.05, 2.55, 0.4, 0.06, WHITE)
rect(s, 9.05, 3.55, 0.4, 0.06, WHITE)

# Bottom: payload + keepalive
rect(s, 0.4, 4.4, 6.2, 2.7, CARD_BG, line=PURPLE, line_w=2)
txt(s, "페이로드 (JSON)", 0.55, 4.5, 5.9, 0.4, size=13, bold=True, color=PURPLE)
multitxt(s, [
    "id · event_id · event_title",
    "success (분석 성공/실패)",
    "stored_at (분석 저장 시각)",
    "read_at (읽음 시각, nullable)",
    "error_message (실패 시)",
], 0.55, 4.95, 5.9, 2.0, size=11, bullet="• ")

rect(s, 6.75, 4.4, 6.2, 2.7, CARD_BG, line=YELLOW, line_w=2)
txt(s, "운영 주의사항", 6.9, 4.5, 5.9, 0.4, size=13, bold=True, color=YELLOW)
multitxt(s, [
    "Keepalive: 30초마다 ': keep-alive\\n\\n'",
    "단일 프로세스 메모리 기반 pub/sub",
    "uvicorn 워커 ≥ 2면 워커 간 공유 ❌",
    "멀티 워커 필요 시 Redis Pub/Sub 등으로 교체",
], 6.9, 4.95, 5.9, 2.0, size=11, bullet="⚠ ")


# ══════════════════════════════════════════════════════════
# Slide 8: 두 기능의 협력 — 매크로 스냅샷 공유
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "두 기능의 협력", "매크로 지표 스냅샷 13종 공유 → 일관된 시장 톤")

# Top: macro feeds
rect(s, 0.4, 1.1, 12.55, 1.7, CARD_BG, line=ACCENT, line_w=2)
rect(s, 0.4, 1.1, 0.25, 1.7, ACCENT)
txt(s, "거시 경제 현황판", 0.7, 1.25, 6.0, 0.4, size=14, bold=True, color=ACCENT)
multitxt(s, [
    "학습 노트 + YouTube + LLM → Risk-on / Risk-off 판단",
    "매크로 지표 13종 스냅샷 (FRED 1순위, Yahoo 폴백)",
    "메모리 + Redis 25h 캐시",
], 0.7, 1.7, 12.0, 1.05, size=11, bullet="• ")

# Arrow
rect(s, 6.55, 2.95, 0.25, 0.4, ACCENT)
txt(s, "▼  공유  ▼", 5.8, 3.0, 1.8, 0.4, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Bottom: notifications consumes
rect(s, 0.4, 3.5, 12.55, 1.7, CARD_BG, line=ORANGE, line_w=2)
rect(s, 0.4, 3.5, 0.25, 1.7, ORANGE)
txt(s, "경제일정 알림 — 영향 분석", 0.7, 3.65, 12.0, 0.4, size=14, bold=True, color=ORANGE)
multitxt(s, [
    "같은 13종 지표를 OpenAIEventImpactAnalyzer 프롬프트에 주입",
    "이 매크로 환경 하에서 일정의 direction · key_drivers · risks 산출",
    "분석 완료 후 Publisher → DB INSERT + SSE Broadcast",
], 0.7, 4.1, 12.0, 1.05, size=11, bullet="• ")

# Indicator chips (shared)
rect(s, 0.4, 5.45, 12.55, 1.65, CARD_BG2, line=GREEN, line_w=2)
txt(s, "공유 매크로 지표 13종", 0.55, 5.55, 12.2, 0.4, size=13, bold=True, color=GREEN)

chips = [
    "DGS10", "US_T2Y", "US_T20Y", "WTI", "GOLD",
    "USD/KRW", "USD/JPY", "DXY", "VIX",
    "S&P 500", "NASDAQ 100", "KOSPI 200", "FRED+Yahoo",
]
chip_left = 0.55
chip_top = 6.0
for i, name in enumerate(chips):
    col = i % 7
    row = i // 7
    x = chip_left + col * 1.78
    y = chip_top + row * 0.5
    rect(s, x, y, 1.65, 0.4, ACCENT)
    txt(s, name, x, y + 0.05, 1.65, 0.3, size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# Slide 9: 운영 메모 + 요약
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
header(s, "운영 메모 & 요약", "Operations · Key Numbers")

# Operations table
rect(s, 0.4, 1.1, 8.0, 4.0, CARD_BG, line=BORDER, line_w=1)
txt(s, "운영 메모", 0.6, 1.2, 7.6, 0.4, size=14, bold=True, color=ACCENT)

rows = [
    ("항목", "값", True),
    ("Macro 갱신 주기", "매일 01:00 KST (APScheduler)", False),
    ("Macro 캐시 TTL", "Redis 25h + 메모리 무제한", False),
    ("Notification 트리거", "LLM 분석 저장 직후", False),
    ("SSE 워커 호환성", "단일 워커 전제", False),
    ("매크로 지표 수", "13종 (FRED 1순위 + Yahoo 폴백)", False),
    ("FOMC 충돌 처리", "(country, date) 1건 collapse", False),
    ("Title 충돌 처리", "윈도우 내 ' (M/D)' suffix", False),
]
top = 1.65
for label, value, is_head in rows:
    bg = HEADER_BG if is_head else CARD_BG2
    rect(s, 0.6, top, 7.6, 0.4, bg)
    txt(s, label, 0.7, top + 0.05, 3.0, 0.3, size=11, bold=is_head, color=WHITE)
    txt(s, value, 3.8, top + 0.05, 4.3, 0.3, size=11, bold=is_head, color=YELLOW if not is_head else WHITE)
    top += 0.42

# Metric cards
rect(s, 8.6, 1.1, 4.4, 4.0, CARD_BG, line=BORDER, line_w=1)
txt(s, "주요 메트릭", 8.8, 1.2, 4.0, 0.4, size=14, bold=True, color=GREEN)

metrics = [
    ("13", "매크로 지표", ACCENT),
    ("25h", "Redis TTL", YELLOW),
    ("4", "API 엔드포인트", ORANGE),
    ("3", "출력 status", PURPLE),
]
for i, (n, label, c) in enumerate(metrics):
    top = 1.8 + i * 0.75
    rect(s, 8.8, top, 1.1, 0.6, c)
    txt(s, n, 8.8, top + 0.1, 1.1, 0.4, size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, label, 10.05, top + 0.15, 2.85, 0.3, size=12, color=WHITE)

# Bottom summary
rect(s, 0.4, 5.3, 12.55, 1.85, CARD_BG2, line=ACCENT, line_w=1)
txt(s, "핵심 요약", 0.6, 5.4, 12.2, 0.4, size=15, bold=True, color=ACCENT)
multitxt(s, [
    "거시 현황판: 일 1회 LLM 갱신 + 25h 캐시로 즉시 응답, contextual + baseline 듀얼 판단",
    "경제일정 알림: 영향 분석 저장 시 DB INSERT + SSE 푸시 (저장과 푸시 분리)",
    "두 기능은 매크로 지표 13종 스냅샷을 공유 → 일관된 시장 톤 + 빠른 분석",
    "출처는 Antelligen AI 자체 분석으로 일원화, 외부 채널/저자명 노출 금지",
], 0.6, 5.8, 12.2, 1.3, size=12, bullet="✓ ")


out = "/Users/a82108/Documents/GitHub/Multi-Agent/antelligen-backend/docs/macro-and-schedule-notifications.pptx"
prs.save(out)
print(f"Saved: {out}")
